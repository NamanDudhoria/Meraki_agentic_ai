import argparse
import os
import re
import uuid
import json
import hashlib
from pathlib import Path
from typing import Any, Iterable, List, Optional, Tuple
from datetime import datetime, timezone

from dotenv import load_dotenv
from langchain_core.documents import Document
from langchain_core.messages import HumanMessage
from langchain_core.tools import tool
from langchain_text_splitters import RecursiveCharacterTextSplitter

from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma

from langchain_anthropic import ChatAnthropic
from langgraph.prebuilt import create_react_agent


WORKSPACE_ROOT = Path(__file__).resolve().parent
load_dotenv(WORKSPACE_ROOT / ".env")

# Your proprietary data
BRAND_DECKS_DIR = WORKSPACE_ROOT / "Brand Solution Decks"
REPORTS_DATA_DIR = WORKSPACE_ROOT / "Reports & Data"

# Local persistent vector store (your "static brain")
VSTORE_DIR = WORKSPACE_ROOT / ".chroma_meraki"


DEFAULT_EMBEDDING_MODEL = os.environ.get("MERAKI_EMBEDDING_MODEL", "sentence-transformers/all-MiniLM-L6-v2")
DEFAULT_CHUNK_SIZE = int(os.environ.get("MERAKI_CHUNK_SIZE", "1000"))
DEFAULT_CHUNK_OVERLAP = int(os.environ.get("MERAKI_CHUNK_OVERLAP", "150"))
DEFAULT_TOP_K = int(os.environ.get("MERAKI_TOP_K", "4"))
VECTOR_BACKEND = os.environ.get("VECTOR_BACKEND", "chroma").strip().lower()
SUPABASE_DB_URL = os.environ.get("SUPABASE_DB_URL", "").strip()
UPSERT_BATCH_SIZE = int(os.environ.get("MERAKI_UPSERT_BATCH_SIZE", "64"))

# OCR tuning (used only when a PDF has little/no extractable text)
MERAKI_OCR_MIN_CHARS = int(os.environ.get("MERAKI_OCR_MIN_CHARS", "800"))
MERAKI_OCR_MAX_PAGES = int(os.environ.get("MERAKI_OCR_MAX_PAGES", "8"))
MERAKI_OCR_DPI = int(os.environ.get("MERAKI_OCR_DPI", "200"))
MERAKI_OCR_LANG = os.environ.get("MERAKI_OCR_LANG", "eng")
MERAKI_ENABLE_OCR = os.environ.get("MERAKI_ENABLE_OCR", "true").strip().lower() in {"1", "true", "yes", "y"}


def _clean_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _iter_files(root: Path, exts: Tuple[str, ...]) -> Iterable[Path]:
    if not root.exists():
        return []
    return (p for p in root.rglob("*") if p.is_file() and p.suffix.lower() in exts)


def _stable_doc_id(source: str, doc_type: str, chunk_index: int) -> str:
    raw = f"{source}|{doc_type}|{chunk_index}"
    return "doc-" + hashlib.sha1(raw.encode("utf-8")).hexdigest()


def _ocr_pdf_text(pdf_path: Path) -> str:
    """
    OCR fallback for scanned PDFs.

    Note: On Windows, `pdf2image` requires Poppler binaries (e.g. `pdftoppm`) to be installed
    and available on PATH.
    """
    from pdf2image import convert_from_path
    import pytesseract

    pages_text: List[str] = []
    images = convert_from_path(
        str(pdf_path),
        dpi=MERAKI_OCR_DPI,
        first_page=1,
        last_page=min(MERAKI_OCR_MAX_PAGES, 10_000),  # guardrail
    )
    for img in images:
        pages_text.append(pytesseract.image_to_string(img, lang=MERAKI_OCR_LANG) or "")
    return _clean_text("\n\n".join(pages_text))


def _should_ocr(extracted_text: str) -> bool:
    if not MERAKI_ENABLE_OCR:
        return False
    extracted_text = (extracted_text or "").strip()
    return len(extracted_text) < MERAKI_OCR_MIN_CHARS


def extract_pdf_text(pdf_path: Path) -> str:
    # Fast path: use pypdf text extraction (works well for "real" PDFs).
    # Fallback: if the PDF is likely scanned (very little extracted text), OCR the first pages.
    from pypdf import PdfReader

    reader = PdfReader(str(pdf_path))
    chunks: List[str] = []
    for page in reader.pages:
        chunks.append(page.extract_text() or "")
    extracted = _clean_text("\n".join(chunks))
    if _should_ocr(extracted):
        try:
            return _ocr_pdf_text(pdf_path)
        except Exception as e:
            # Don't hard-fail ingestion if OCR dependencies aren't installed.
            # Keep whatever text we could extract.
            print(f"[OCR fallback failed] {pdf_path.name}: {type(e).__name__}: {e}")
    return extracted


def extract_docx_text(docx_path: Path) -> str:
    from docx import Document as DocxDocument

    doc = DocxDocument(str(docx_path))
    paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    return _clean_text("\n".join(paragraphs))


def extract_pptx_text(pptx_path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slide_texts: List[str] = []

    for slide in prs.slides:
        parts: List[str] = []
        for shape in slide.shapes:
            text = ""
            # Some shapes expose .text directly; text frames expose shape.text too.
            try:
                text = shape.text or ""
            except Exception:
                text = ""
            if text and text.strip():
                parts.append(text.strip())
        if parts:
            slide_texts.append("\n".join(parts))

    return _clean_text("\n\n".join(slide_texts))


def load_documents(errors: Optional[List[str]] = None) -> List[Document]:
    documents: List[Document] = []
    _errors = errors if errors is not None else []

    pdfs = list(_iter_files(BRAND_DECKS_DIR, (".pdf",))) + list(_iter_files(REPORTS_DATA_DIR, (".pdf",)))
    docxs = list(_iter_files(BRAND_DECKS_DIR, (".docx",))) + list(_iter_files(REPORTS_DATA_DIR, (".docx",)))
    pptxs = list(_iter_files(BRAND_DECKS_DIR, (".pptx",))) + list(_iter_files(REPORTS_DATA_DIR, (".pptx",)))

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=DEFAULT_CHUNK_SIZE,
        chunk_overlap=DEFAULT_CHUNK_OVERLAP,
        separators=["\n\n", "\n", " ", ""],
    )

    def add_text(text: str, *, source: str, doc_type: str) -> None:
        if not text or len(text) < 50:
            return
        for idx, chunk in enumerate(splitter.split_text(text)):
            chunk = _clean_text(chunk)
            if chunk:
                documents.append(
                    Document(
                        page_content=chunk,
                        metadata={
                            "source": source,
                            "type": doc_type,
                            "chunk_index": idx,
                        },
                    )
                )

    for p in pdfs:
        try:
            text = extract_pdf_text(p)
            add_text(text, source=str(p), doc_type="pdf")
        except Exception as e:
            msg = f"[PDF] {p.name}: {type(e).__name__}: {e}"
            print(f"[PDF extraction failed] {msg}")
            _errors.append(msg)

    for p in docxs:
        try:
            text = extract_docx_text(p)
            add_text(text, source=str(p), doc_type="docx")
        except Exception as e:
            msg = f"[DOCX] {p.name}: {type(e).__name__}: {e}"
            print(f"[DOCX extraction failed] {msg}")
            _errors.append(msg)

    for p in pptxs:
        try:
            text = extract_pptx_text(p)
            add_text(text, source=str(p), doc_type="pptx")
        except Exception as e:
            msg = f"[PPTX] {p.name}: {type(e).__name__}: {e}"
            print(f"[PPTX extraction failed] {msg}")
            _errors.append(msg)

    return documents


class PostgresVectorStore:
    def __init__(self, db_url: str, embeddings: HuggingFaceEmbeddings):
        self.db_url = db_url
        self.embeddings = embeddings

    def _connect(self):
        import psycopg

        return psycopg.connect(self.db_url)

    @staticmethod
    def _as_vector_literal(values: List[float]) -> str:
        return "[" + ",".join(f"{float(v):.8f}" for v in values) + "]"

    def similarity_search(self, query: str, k: int = 4, filter: Optional[dict] = None) -> List[Document]:
        qvec = self.embeddings.embed_query(query)
        vector_literal = self._as_vector_literal(qvec)
        filt = filter or {}
        docs: List[Document] = []
        with self._connect() as conn:
            with conn.cursor() as cur:
                cur.execute(
                    """
                    select id, content, metadata, similarity
                    from match_meraki_documents(%s::vector, %s, %s::jsonb)
                    """,
                    (vector_literal, int(k), json.dumps(filt)),
                )
                rows = cur.fetchall()
        for row in rows:
            _, content, metadata, _ = row
            docs.append(Document(page_content=content, metadata=metadata or {}))
        return docs

    def add_documents(self, docs: List[Document], ids: Optional[List[str]] = None) -> None:
        if not docs:
            return
        texts = [d.page_content for d in docs]
        vectors = self.embeddings.embed_documents(texts)
        resolved_ids = ids or [
            _stable_doc_id(
                source=str(d.metadata.get("source", "unknown")),
                doc_type=str(d.metadata.get("type", "unknown")),
                chunk_index=int(d.metadata.get("chunk_index", i)),
            )
            for i, d in enumerate(docs)
        ]
        sql = """
            insert into meraki_documents (id, content, metadata, embedding, updated_at)
            values (%s, %s, %s::jsonb, %s::vector, now())
            on conflict (id) do update set
              content = excluded.content,
              metadata = excluded.metadata,
              embedding = excluded.embedding,
              updated_at = now()
        """
        params = [
            (
                doc_id,
                doc.page_content,
                json.dumps(doc.metadata or {}),
                self._as_vector_literal(vec),
            )
            for doc, doc_id, vec in zip(docs, resolved_ids, vectors)
        ]
        total = len(params)
        with self._connect() as conn:
            with conn.cursor() as cur:
                for i in range(0, total, max(1, UPSERT_BATCH_SIZE)):
                    batch = params[i : i + max(1, UPSERT_BATCH_SIZE)]
                    cur.executemany(sql, batch)
                    conn.commit()
                    print(f"[Supabase upsert] {min(i + len(batch), total)}/{total}")

    def delete(self, ids: List[str]) -> None:
        if not ids:
            return
        with self._connect() as conn:
            with conn.cursor() as cur:
                cur.execute("delete from meraki_documents where id = any(%s)", (ids,))
            conn.commit()

    def clear_all(self) -> None:
        with self._connect() as conn:
            with conn.cursor() as cur:
                cur.execute("delete from meraki_documents")
            conn.commit()


def _diagnose_empty(extraction_errors: List[str]) -> str:
    brand_files = list(_iter_files(BRAND_DECKS_DIR, (".pdf", ".docx", ".pptx")))
    reports_files = list(_iter_files(REPORTS_DATA_DIR, (".pdf", ".docx", ".pptx")))
    total_files = len(brand_files) + len(reports_files)

    lines = ["No text could be extracted from your data folders."]
    lines.append(
        f"  Brand Solution Decks/ — exists={BRAND_DECKS_DIR.exists()}, {len(brand_files)} file(s)"
    )
    lines.append(
        f"  Reports & Data/       — exists={REPORTS_DATA_DIR.exists()}, {len(reports_files)} file(s)"
    )

    if total_files == 0:
        lines.append(
            "No supported files (.pdf, .docx, .pptx) found. "
            "Make sure these folders are present in your deployment."
        )
    elif extraction_errors:
        lines.append(f"{len(extraction_errors)} extraction error(s) (first 10 shown):")
        for e in extraction_errors[:10]:
            lines.append(f"  • {e}")
    else:
        lines.append(
            f"All {total_files} file(s) produced fewer than 50 characters of text — "
            "they are likely scanned images. "
            "OCR requires Poppler (pdftoppm) + Tesseract to be installed."
        )
    return "\n".join(lines)


def build_vectorstore(rebuild: bool = False) -> Any:
    embeddings = HuggingFaceEmbeddings(model_name=DEFAULT_EMBEDDING_MODEL)

    if VECTOR_BACKEND in {"postgres", "supabase"}:
        if not SUPABASE_DB_URL:
            raise RuntimeError("VECTOR_BACKEND=postgres requires SUPABASE_DB_URL in environment.")
        vs = PostgresVectorStore(SUPABASE_DB_URL, embeddings)
        if rebuild:
            vs.clear_all()
            extraction_errors: List[str] = []
            documents = load_documents(errors=extraction_errors)
            if not documents:
                raise RuntimeError(_diagnose_empty(extraction_errors))
            ids = [
                _stable_doc_id(
                    source=str(d.metadata.get("source", "unknown")),
                    doc_type=str(d.metadata.get("type", "unknown")),
                    chunk_index=int(d.metadata.get("chunk_index", i)),
                )
                for i, d in enumerate(documents)
            ]
            # Upsert in chunks so long ingests persist incremental progress.
            chunk = max(1, UPSERT_BATCH_SIZE)
            total = len(documents)
            for i in range(0, total, chunk):
                docs_batch = documents[i : i + chunk]
                ids_batch = ids[i : i + chunk]
                vs.add_documents(docs_batch, ids=ids_batch)
                print(f"[Rebuild progress] {min(i + len(docs_batch), total)}/{total} docs")
            print("Ingestion complete. Rebuild index ready to use with --query or interactive mode.")
        return vs

    if rebuild and VSTORE_DIR.exists():
        for child in VSTORE_DIR.glob("**/*"):
            if child.is_file():
                child.unlink(missing_ok=True)
        for child in sorted(VSTORE_DIR.glob("**/*"), reverse=True):
            if child.is_dir():
                try:
                    child.rmdir()
                except OSError:
                    pass

    if VSTORE_DIR.exists() and any(VSTORE_DIR.iterdir()) and not rebuild:
        return Chroma(persist_directory=str(VSTORE_DIR), embedding_function=embeddings, collection_name="meraki_history")

    extraction_errors = []
    documents = load_documents(errors=extraction_errors)
    if not documents:
        raise RuntimeError(_diagnose_empty(extraction_errors))

    VSTORE_DIR.mkdir(parents=True, exist_ok=True)
    vs = Chroma.from_documents(
        documents=documents,
        embedding=embeddings,
        persist_directory=str(VSTORE_DIR),
        collection_name="meraki_history",
    )
    # persist() was removed in chromadb ≥ 0.4; persistence is now automatic.
    if hasattr(vs, "persist"):
        vs.persist()
    return vs


def create_rag_agent(
    vectorstore: Any,
    session_id: Optional[str] = None,
    include_all_memory: bool = True,
):
    llm_model = os.environ.get("ANTHROPIC_MODEL", "claude-sonnet-4-6")
    llm = ChatAnthropic(model=llm_model, temperature=0)

    @tool
    def meraki_history_search(query: str) -> str:
        """
        Search your Meraki historical solutions (brand decks, reports, and prior work) and return the most relevant excerpts.
        Always use this when you need to reference Meraki's past methodology, campaign patterns, funnels, or examples.
        """
        try:
            docs = vectorstore.similarity_search(query, k=max(DEFAULT_TOP_K * 3, 12))
        except Exception as e:
            return f"[Database error — could not search Meraki knowledge base: {type(e).__name__}: {e}]"
        if not include_all_memory and session_id:
            filtered_docs: List[Document] = []
            for d in docs:
                if d.metadata.get("type") != "chat_memory":
                    filtered_docs.append(d)
                elif d.metadata.get("session_id") == session_id:
                    filtered_docs.append(d)
            docs = filtered_docs
        docs = docs[:DEFAULT_TOP_K]
        if not docs:
            return "No relevant Meraki historical excerpts found."

        formatted: List[str] = []
        for i, d in enumerate(docs, start=1):
            source = d.metadata.get("source", "unknown")
            doc_type = d.metadata.get("type", "unknown")
            excerpt = d.page_content
            if len(excerpt) > 1200:
                excerpt = excerpt[:1200].rstrip() + "..."
            formatted.append(f"[{i}] ({doc_type}) {source}\n{excerpt}")
        return "\n\n".join(formatted)

    @tool
    def web_search(query: str) -> str:
        """
        Search the live web for current sports/business/news/trends.

        Use this for competitor research, "today's" updates, recent partnerships, and fast market context.
        Returns top snippets with URLs so you can cite sources in your answer.
        """
        results: List[str] = []

        # Tavily (primary)
        tavily_key = os.environ.get("TAVILY_API_KEY", "").strip()
        if tavily_key:
            try:
                from tavily import TavilyClient
                client = TavilyClient(api_key=tavily_key)
                response = client.search(query, max_results=6)
                for i, r in enumerate(response.get("results", []), start=1):
                    title = r.get("title") or "Untitled"
                    url = r.get("url") or ""
                    content = _clean_text(r.get("content") or r.get("snippet") or "")[:500]
                    results.append(f"[{i}] {title}\nURL: {url}\n{content}")
                if results:
                    return "\n\n".join(results)
            except Exception:
                pass  # fall through to DuckDuckGo

        # DuckDuckGo (fallback)
        try:
            from duckduckgo_search import DDGS
            with DDGS() as ddgs:
                first_batch = list(ddgs.text(query, max_results=6) or [])
                if not first_batch:
                    first_batch = list(ddgs.news(query, max_results=6) or [])
                for i, r in enumerate(first_batch, start=1):
                    if not isinstance(r, dict):
                        continue
                    title = r.get("title") or "Untitled"
                    href = r.get("href") or r.get("url") or ""
                    body = _clean_text(r.get("body") or r.get("description") or r.get("snippet") or "")[:500]
                    results.append(f"[{i}] {title}\nURL: {href}\n{body}")
        except Exception as e:
            return f"[Web search unavailable: {type(e).__name__}. Answering from Meraki knowledge base only.]"

        return "\n\n".join(results) if results else "No web results found."

    tools = [meraki_history_search, web_search]

    system_prompt = (
        "You are an elite sports marketing consultant modeled after Meraki Sport & Entertainment.\n"
        "Your expertise is identifying brand-sport partnerships, conceptualizing fan engagement funnels,\n"
        "and optimizing team social platforms.\n\n"
        "Use the tool `meraki_history_search` whenever you need Meraki-specific historical examples or strategy patterns.\n"
        "Use the tool `web_search` when you need live competitor data, current sports/business news, or today's trends.\n"
        "Be concise, actionable, and always ground recommendations in the retrieved excerpts when possible."
    )

    # LangGraph ReAct loop: the LLM decides when to call tools, then synthesizes the final answer.
    agent = create_react_agent(llm, tools, prompt=system_prompt)
    return agent


def remember_chat_interaction(
    vectorstore: Any,
    user_message: str,
    assistant_message: str,
    session_id: str = "default",
) -> Optional[str]:
    """
    Persist each chat turn into the same vector store so future queries can retrieve it.
    """
    user_message = _clean_text(user_message or "")
    assistant_message = _clean_text(assistant_message or "")
    if not user_message and not assistant_message:
        return None

    memory_id = f"chat-{session_id}-{uuid.uuid4()}"
    memory_text = (
        "Chat Memory\n"
        f"Session: {session_id}\n"
        f"Timestamp: {datetime.now(timezone.utc).isoformat()}\n\n"
        f"User:\n{user_message}\n\n"
        f"Assistant:\n{assistant_message}"
    )
    memory_doc = Document(
        page_content=memory_text,
        metadata={
            "type": "chat_memory",
            "source": f"chat:{session_id}",
            "session_id": session_id,
            "memory_id": memory_id,
        },
    )
    vectorstore.add_documents([memory_doc], ids=[memory_id])
    return memory_id


def forget_memory_record(vectorstore: Any, memory_id: str) -> None:
    """
    Remove a previously stored chat-memory record by id.
    """
    if not memory_id:
        return
    vectorstore.delete(ids=[memory_id])


def _parse_memory_text(text: str) -> dict:
    """Parse the stored memory_text format back into structured fields."""
    result = {"session_id": "", "timestamp": "", "user": "", "assistant": ""}
    for line in text.split("\n")[:5]:
        if line.startswith("Session: "):
            result["session_id"] = line[9:].strip()
        elif line.startswith("Timestamp: "):
            result["timestamp"] = line[11:].strip()
    user_marker = "\nUser:\n"
    asst_marker = "\nAssistant:\n"
    u = text.find(user_marker)
    a = text.find(asst_marker)
    if u != -1:
        result["user"] = text[u + len(user_marker): a if a != -1 else None].strip()
    if a != -1:
        result["assistant"] = text[a + len(asst_marker):].strip()
    return result


def get_chat_sessions(vectorstore: Any) -> List[dict]:
    """
    Return all saved chat sessions stored in the vector store, sorted newest-first.
    Each entry: {"session_id", "messages": [{"role", "content"}], "preview", "timestamp"}
    """
    from collections import defaultdict

    raw: List[dict] = []

    if VECTOR_BACKEND in {"postgres", "supabase"}:
        try:
            with vectorstore._connect() as conn:
                with conn.cursor() as cur:
                    cur.execute(
                        "SELECT content, metadata FROM meraki_documents "
                        "WHERE metadata->>'type' = 'chat_memory' ORDER BY updated_at ASC"
                    )
                    for content, metadata in cur.fetchall():
                        raw.append({"content": content or "", "metadata": metadata or {}})
        except Exception as e:
            print(f"[get_chat_sessions] Postgres error: {e}")
    else:
        try:
            result = vectorstore._collection.get(
                where={"type": {"$eq": "chat_memory"}},
                include=["documents", "metadatas"],
            )
            for content, meta in zip(result.get("documents") or [], result.get("metadatas") or []):
                raw.append({"content": content or "", "metadata": meta or {}})
        except Exception as e:
            print(f"[get_chat_sessions] Chroma error: {e}")

    sessions: dict = defaultdict(lambda: {"messages": [], "latest_ts": "", "preview": ""})
    for item in raw:
        parsed = _parse_memory_text(item["content"])
        sid = parsed["session_id"] or item["metadata"].get("session_id", "unknown")
        if parsed["user"]:
            sessions[sid]["messages"].append({"role": "user", "content": parsed["user"]})
        if parsed["assistant"]:
            sessions[sid]["messages"].append({"role": "assistant", "content": parsed["assistant"]})
        if parsed["timestamp"] > sessions[sid]["latest_ts"]:
            sessions[sid]["latest_ts"] = parsed["timestamp"]
        if not sessions[sid]["preview"] and parsed["user"]:
            sessions[sid]["preview"] = parsed["user"][:80]

    out = [
        {
            "session_id": sid,
            "messages": data["messages"],
            "preview": data["preview"],
            "timestamp": data["latest_ts"],
        }
        for sid, data in sessions.items()
    ]
    out.sort(key=lambda x: x["timestamp"], reverse=True)
    return out


def run_interactive(agent):
    print("Meraki RAG agent ready. Type a question (or 'exit').")
    while True:
        q = input("\nYou: ").strip()
        if not q:
            continue
        if q.lower() in {"exit", "quit"}:
            return
        result = agent.invoke({"messages": [HumanMessage(content=q)]})
        final_msg = result["messages"][-1]
        print("\nAgent:", getattr(final_msg, "content", str(final_msg)).strip())


def main():
    parser = argparse.ArgumentParser(description="Meraki Agentic RAG (static brain from your decks/reports).")
    parser.add_argument("--ingest", action="store_true", help="Build/rebuild the local vector index.")
    parser.add_argument("--rebuild", action="store_true", help="Rebuild the index from scratch.")
    parser.add_argument("--query", type=str, default=None, help="Single question to ask the agent.")
    args = parser.parse_args()

    vectorstore = build_vectorstore(rebuild=args.rebuild or args.ingest)

    # If the user only asked to ingest/rebuild, do not require an LLM key and do not start interactive mode.
    if (args.ingest or args.rebuild) and args.query is None:
        print("Ingestion complete. Rebuild index ready to use with --query or interactive mode.")
        return

    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        raise RuntimeError(
            "Missing ANTHROPIC_API_KEY. Set it in your environment before running.\n"
            "Example (PowerShell): $env:ANTHROPIC_API_KEY='your_key'"
        )

    agent = create_rag_agent(vectorstore)
    if args.query:
        result = agent.invoke({"messages": [HumanMessage(content=args.query)]})
        final_msg = result["messages"][-1]
        print(getattr(final_msg, "content", str(final_msg)).strip())
    else:
        run_interactive(agent)


if __name__ == "__main__":
    main()
